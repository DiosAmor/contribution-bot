# teleconference-agenda, meeting-agenda를 기반으로 하여 기고문 분류
# 이 때, 파일 형식 (pptx, docx)에 따라서 다르게 활용.
# 다운로드 링크는 crawling excel 파일에서 검색하여 활용.
import openpyxl
from docx import Document
from openpyxl.styles import Font
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def edit_excel(filename, sheetname, update_list):
    header_font = Font(bold=True, size=14)
    excel_file = openpyxl.load_workbook(filename)
    excel_sheet = excel_file[sheetname]

    excel_sheet.cell(row=1, column=11, value="Topic").font = header_font
    excel_sheet.cell(row=1, column=12, value="Session").font = header_font
    for row in range(2, excel_sheet.max_row + 1):
        year = excel_sheet.cell(row=row, column=2).value
        dcn = excel_sheet.cell(row=row, column=3).value
        for update_item in update_list:
            if int(update_item[0]) == int(year) and int(update_item[1]) == int(dcn):
                excel_sheet.cell(row=row, column=11, value=update_item[2])
                excel_sheet.cell(row=row, column=12, value=update_item[3])

    excel_file.save(filename="test.xlsx")


def open_teleconference_agenda(filename):
    doc = Document(filename)
    submission_list = []

    # 세 번째 테이블이 기고문 리스트
    table = doc.tables[2]
    for row in table.rows:
        year = row.cells[0].text.split("/")[0]
        if not year.isdigit():
            continue
        if len(year) > 2:
            dcn = year
            year = "2023"
        else:
            dcn = row.cells[0].text.split("/")[1]
            year = "20" + year
        topic = row.cells[4].text
        session = row.cells[5].text
        items = [year, dcn, topic, session]
        submission_list.append(items)
    return submission_list


def open_meeting_agenda(filename):
    prs = Presentation(filename)
    submission_list = []
    # 슬라이스 숫자 0부터 셈
    # 후에는 submission 글자가 포함되어있는 슬라이드부터 세는 것도 나쁘지 않을듯.
    for submission_idx in range(19, 39):
        slide = prs.slides[submission_idx]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                for row in table.rows:
                    year = row.cells[0].text.split("/")[0]
                    if not year.isdigit():
                        continue
                    if len(year) > 2:
                        dcn = year
                        year = "2023"
                    else:
                        dcn = row.cells[0].text.split("/")[1]
                        year = "20" + year
                    topic = row.cells[4].text
                    session = row.cells[5].text
                    items = [year, dcn, topic, session]
                    submission_list.append(items)
    return submission_list


def do_arrangement():
    filename_list = [
        "11-23-1713-14-00bn-tgbn-nov-2023-meeting-agenda.pptx",
        "11-23-2140-10-00bn-nov-jan-tgbn-teleconference-agenda.docx",
        "11-23-2174-11-00bn-tgbn-jan-2024-meeting-agenda.pptx",
        "11-24-0201-11-00bn-jan-mar-tgbn-teleconference-agenda.docx",
        "11-24-0235-14-00bn-tgbn-mar-2024-meeting-agenda.pptx",
        "11-24-0633-15-00bn-mar-may-tgbn-teleconference-agenda.docx",
        "11-24-0653-15-00bn-tgbn-may-2024-meeting-agenda.pptx",
        "11-24-0964-05-00bn-may-july-tgbn-teleconference-agenda.docx",
    ]
    excel_filename = "test.xlsx"
    for filename in filename_list:
        print(filename)
        if filename.endswith(".docx"):
            submission_list = open_teleconference_agenda("agenda/" + filename)
        else:
            submission_list = open_meeting_agenda("agenda/" + filename)
        edit_excel(excel_filename, "00bn", submission_list)


do_arrangement()
