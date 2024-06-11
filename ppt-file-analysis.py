from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# PowerPoint 파일 열기
prs = Presentation("11-24-0653-15-00bn-tgbn-may-2024-meeting-agenda.pptx")

# 모든 슬라이드에서 텍스트 읽기

# Slicing이 안 되는데 뭐지?
# getitem 쪽은 더 파봐야하는듯? rId가 뭐여
# 일단은 그냥 쓰자
submission_list = []
for submission_idx in range(19, 31):
    slide = prs.slides[submission_idx]
    print(slide)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                items = []
                for cell in row.cells:
                    items.append(cell.text)
                submission_list.append(items)
print(submission_list)
